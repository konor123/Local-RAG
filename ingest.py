import faulthandler
faulthandler.enable()
import os
import sys

# Disable ChromaDB Telemetry via environment variables
os.environ["ANONYMIZED_TELEMETRY"] = "False"
os.environ["CHROMA_SERVER_NOINTERACTIVE"] = "True"
os.environ["TOKENIZERS_PARALLELISM"] = "false"
os.environ["OMP_NUM_THREADS"] = "4"  # Ryzen 5 5600G optimization

# Fix Windows console UTF-8 mojibake (safe for Python 3.7+)
if sys.platform == "win32" and hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")

import glob
from typing import List, Dict, Any
from drive_manager import filter_walk_dirs, get_exclude_dir_names, get_search_roots
from runtime_paths import logs_dir, runtime_path
from langchain_community.document_loaders import (
    PyPDFLoader,
    Docx2txtLoader,
    UnstructuredExcelLoader,
    UnstructuredHTMLLoader,
    TextLoader,
)
from langchain_community.document_loaders.base import BaseLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_core.documents import Document
import re
import subprocess
import sys
import shutil
import tempfile
import os

# --- STABILITY PATCH ---
# Disable tokenizer parallelism to prevent crashes in multi-threaded/process environments
os.environ["TOKENIZERS_PARALLELISM"] = "false"
# Limit BLAS/Torch threads to avoid contention
os.environ.setdefault("OMP_NUM_THREADS", "1")
os.environ.setdefault("MKL_NUM_THREADS", "1")
# -----------------------

from concurrent.futures import ThreadPoolExecutor, ProcessPoolExecutor, as_completed, wait, FIRST_COMPLETED

# --- Global: Embedder Process Logic (Must be top-level for Pickle) ---
_process_embed_model = None

def init_embedder_process():
    """Initializer for Embedder Process"""
    global _process_embed_model
    try:
        # Suppress parallelism in worker
        os.environ["TOKENIZERS_PARALLELISM"] = "false"
        os.environ["OMP_NUM_THREADS"] = "1"
        os.environ["ANONYMIZED_TELEMETRY"] = "False"
        os.environ["CHROMA_SERVER_NOINTERACTIVE"] = "True"
        
        from langchain_community.embeddings import HuggingFaceEmbeddings
        _process_embed_model = HuggingFaceEmbeddings(
            model_name="dragonkue/multilingual-e5-small-ko",
            model_kwargs={'device': 'cpu'},
            encode_kwargs={'normalize_embeddings': True}
        )
        print(f"   [Embedder-Process] Model Initialized.")
    except Exception as e:
        print(f"   [Embedder-Process] Init Failed: {e}")

def run_embedding_task(texts):
    """Executes embedding in isolated process"""
    global _process_embed_model
    if not _process_embed_model:
        raise RuntimeError("Model not initialized in process")
        
    start = time.time()
    embeddings = _process_embed_model.embed_documents(texts)
    duration = time.time() - start
    return embeddings, duration
# -------------------------------------------------------------------

# ... (omitted lines)

from pptx import Presentation
import json
import logging
import contextlib
import itertools
import queue
import threading
import time
import uuid
from collections import defaultdict, Counter
import csv

# Suppress noisy loggers
logging.getLogger("unstructured").setLevel(logging.ERROR)
logging.getLogger("pdfminer").setLevel(logging.ERROR)
logging.getLogger("PIL").setLevel(logging.ERROR)
logging.getLogger("chromadb").setLevel(logging.ERROR)

# Context manager to suppress stdout/stderr from C-libraries or noisy libs
@contextlib.contextmanager
def suppress_output():
    with open(os.devnull, "w") as devnull:
        old_stdout = sys.stdout
        old_stderr = sys.stderr
        try:
            sys.stdout = devnull
            sys.stderr = devnull
            yield
        finally:
            sys.stdout = old_stdout
            sys.stderr = old_stderr

def _find_libreoffice() -> str:
    r"""Locate the LibreOffice `program` directory in this priority order.

    1. Environment variable `OSL_LIBREOFFICE_PATH` (full path to `program` dir).
    2. Windows registry (HKLM\SOFTWARE\LibreOffice\LibreOffice\...InstallLocation).
    3. Common install locations on Windows.
    4. Fallback: search PATH for `soffice.exe`.
    """
    # 1. Environment variable
    env_path = os.environ.get("OSL_LIBREOFFICE_PATH", "").strip()
    if env_path and os.path.isdir(env_path):
        return env_path

    # 2. Windows registry
    if sys.platform == "win32":
        try:
            import winreg
            for sub_key in (
                (winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\LibreOffice\LibreOffice"),
                (winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\WOW6432Node\LibreOffice\LibreOffice"),
            ):
                try:
                    with winreg.OpenKey(*sub_key) as key:
                        install_path, _ = winreg.QueryValueEx(key, "InstallLocation")
                    if install_path:
                        program = os.path.join(install_path, "program")
                        if os.path.isdir(program):
                            return program
                except OSError:
                    continue
        except Exception:
            pass

    # 3. Common Windows install locations
    if sys.platform == "win32":
        for candidate in (
            r"C:\Program Files\LibreOffice\program",
            r"C:\Program Files (x86)\LibreOffice\program",
        ):
            if os.path.isdir(candidate):
                return candidate

    # 4. PATH search for soffice.exe
    for exe_name in ("soffice.exe", "soffice"):
        for path_dir in os.environ.get("PATH", "").split(os.pathsep):
            if not path_dir:
                continue
            candidate = os.path.join(path_dir, exe_name)
            if os.path.isfile(candidate):
                return path_dir
    return ""


# Add LibreOffice to PATH if found
libreoffice_path = _find_libreoffice()
if libreoffice_path and libreoffice_path not in os.environ["PATH"]:
    print(f"Adding LibreOffice to PATH: {libreoffice_path}")
    os.environ["PATH"] += ";" + libreoffice_path

# Override print to handle Unicode errors in Windows Console
_original_print = print
def safe_print(*args, **kwargs):
    try:
        _original_print(*args, **kwargs)
    except Exception:
        try:
            # Fallback: Encode to UTF-8 then decode ignoring errors, keeps it safe
            clean_args = []
            for arg in args:
                if isinstance(arg, str):
                    clean_args.append(arg.encode('utf-8', 'replace').decode('utf-8'))
                else:
                    clean_args.append(arg)
            _original_print(*clean_args, **kwargs)
        except Exception:
            pass # Silence if even that fails

print = safe_print

def send_windows_notification(title, message):
    """Sends a Windows Toast Notification using PowerShell"""
    try:
        ps_script = f"""
        [Windows.UI.Notifications.ToastNotificationManager, Windows.UI.Notifications, ContentType = WindowsRuntime] > $null
        $template = [Windows.UI.Notifications.ToastNotificationManager]::GetTemplateContent([Windows.UI.Notifications.ToastTemplateType]::ToastText02)
        $textNodes = $template.GetElementsByTagName("text")
        $textNodes.Item(0).AppendChild($template.CreateTextNode("{title}")) > $null
        $textNodes.Item(1).AppendChild($template.CreateTextNode("{message}")) > $null
        $notifier = [Windows.UI.Notifications.ToastNotificationManager]::CreateToastNotifier("OSL AI Assistant Ingest")
        $notifier.Show([Windows.UI.Notifications.ToastNotification($template)])
        """
        subprocess.run(["powershell", "-Command", ps_script], check=False,
                      creationflags=subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0)
    except Exception as e:
        print(f"Failed to send notification: {e}")

# --- Sanitization for Stability ---
def verify_and_sanitize_content(text: str) -> str:
    """
    Fundamentally prevents crashes by removing 'Toxic' patterns 
    that cause Tokenizer/Model segfaults.
    """
    if not text: return ""
    
    # 1. Remove Null Bytes (Common C-string terminator exploit/bug)
    text = text.replace("\x00", "")
    
    # 2. Limit Total Length (Safety cap: 100k chars per 'Document' before split)
    # Most valid docs are split, but if a single 'page' is massive, it kills.
    if len(text) > 100000:
        text = text[:100000] + " [TRUNCATED_SAFETY]"
        
    # 3. Break Massive Words (Binary strings in text disguise)
    # If a single word is > 1000 chars, it's likely garbage. Tokenizer chokes on it.
    # Simple heuristic: inject space every 500 chars if no space found?
    # Doing this via regex is slow. Just truncate very long chunks in splitting?
    # Better: RecursiveCharacterTextSplitter handles this IF configured well.
    # We will enforce it in the worker logic.
    
    return text

class DWGLoader(BaseLoader):
    """
    A simple loader for DWG files that extracts printable strings from the binary content.
    """
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Document]:
        try:
            with open(self.file_path, 'rb') as f:
                content = f.read()
                text = ""
                try:
                    decoded = content.decode('utf-8', errors='ignore')
                    matches = re.findall(r'[\x20-\x7E\uAC00-\uD7A3]{4,}', decoded)
                    text = "\n".join(matches)
                except Exception:
                    pass
                
                if not text:
                    text = "[DWG File - No text extracted]"
                    
                metadata = {"source": self.file_path, "file_type": "dwg"}
                return [Document(page_content=text, metadata=metadata)]
        except Exception as e:
            print(f"Error loading DWG {self.file_path}: {e}")
            return []

class PPTXLoader(BaseLoader):
    """
    A lightweight loader for .pptx files using python-pptx.
    """
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Document]:
        try:
            prs = Presentation(self.file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            
            full_text = "\n".join(text)
            metadata = {"source": self.file_path, "file_type": "pptx"}
            return [Document(page_content=full_text, metadata=metadata)]
        except Exception as e:
            print(f"Error loading PPTX {self.file_path}: {e}")
            return []

def convert_with_soffice(file_path: str, format="txt") -> str:
    """
    Converts a file to text using LibreOffice headless mode.
    Returns the extracted text.
    Handles encoding robustly.
    """
    # Create a temporary directory
    with tempfile.TemporaryDirectory() as temp_dir:
        # Run soffice conversion
        # Syntax: soffice --headless --convert-to txt:Text --outdir <DIR> <FILE>
        cmd = [
            "soffice",
            "--headless",
            "--convert-to", "txt:Text",
            "--outdir", temp_dir,
            file_path
        ]
        
        try:
            # Capture output to avoid console spam, but we typically ignore it unless error
            # Add timeout to prevent hanging on corrupt files
            subprocess.run(
                cmd,
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=60,
                creationflags=subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0,
            )
            
            # Find the result file
            # Original: My File.ppt -> My File.txt
            base_name = os.path.splitext(os.path.basename(file_path))[0]
            # Handle potential name mangling by soffice? 
            # Usually it keeps the name. safely glob it.
            result_files = glob.glob(os.path.join(temp_dir, "*.txt"))
            
            if not result_files:
                raise Exception("Conversion failed, no output file created.")
                
            txt_file = result_files[0]
            
            # Read with robust encoding detection
            content = ""
            try:
                with open(txt_file, "r", encoding="utf-8") as f:
                    content = f.read()
            except UnicodeDecodeError:
                try:
                    with open(txt_file, "r", encoding="cp949") as f:
                        content = f.read()
                except Exception:
                    # Fallback
                    with open(txt_file, "r", encoding="utf-8", errors="replace") as f:
                        content = f.read()
            return content
            
        except subprocess.CalledProcessError as e:
            # print(f"Soffice failed for {file_path}: {e}")
            return ""
        except Exception as e:
            # print(f"Conversion error {file_path}: {e}")
            return ""

class LegacyPPTLoader(BaseLoader):
     """
     Fallback for .ppt files using direct LibreOffice conversion.
     """
     def __init__(self, file_path: str):
         self.file_path = file_path
         
     def load(self) -> List[Document]:
         text = convert_with_soffice(self.file_path)
         if text.strip():
             metadata = {"source": self.file_path, "file_type": "ppt"}
             return [Document(page_content=text, metadata=metadata)]
         return []

class LegacyDocLoader(BaseLoader):
     """
     Fallback for .doc files using direct LibreOffice conversion.
     """
     def __init__(self, file_path: str):
         self.file_path = file_path
         
     def load(self) -> List[Document]:
         text = convert_with_soffice(self.file_path)
         if text.strip():
             metadata = {"source": self.file_path, "file_type": "doc"}
             return [Document(page_content=text, metadata=metadata)]
         return []

class HwpLoader(BaseLoader):
     """Loader for HWP/HWPX files using hwpkit."""
     def __init__(self, file_path: str):
         self.file_path = file_path

     def load(self) -> List[Document]:
         try:
             from hwpkit import extract_text_from_file

             text = extract_text_from_file(self.file_path)
             if isinstance(text, (list, tuple)):
                 text = "\n".join(str(item) for item in text if item)
             text = str(text or "")
             if text.strip():
                 ext = os.path.splitext(self.file_path)[1].lower().lstrip(".")
                 metadata = {"source": self.file_path, "file_type": ext or "hwp"}
                 return [Document(page_content=text, metadata=metadata)]
             return []
         except Exception as e:
             print(f"Error loading HWP/HWPX {self.file_path}: {e}")
             return []

import traceback

def send_status_notification(message):
    """Remote notifications are disabled for the internal local build."""
    return

# Configuration
DRIVES = get_search_roots()
VECTOR_STORE_PATH = os.environ.get("VECTOR_STORE_PATH", runtime_path("chroma_db_ko"))  # Korean-optimized embeddings
TEST_LIMIT = int(os.environ.get("INDEX_SAMPLE_LIMIT", "0") or 0)  # 0 means full indexing
INDEX_SAMPLE_MODE = os.environ.get("INDEX_SAMPLE_MODE", "priority_first").lower()
# Sample runs are report-only by default so processed tracker and vector index are not polluted.
INDEX_SAMPLE_WRITE = os.environ.get("INDEX_SAMPLE_WRITE", "0" if TEST_LIMIT else "1").lower() in ("1", "true", "yes", "y")
WORKER_FILE_TIMEOUT_SEC = int(os.environ.get("WORKER_FILE_TIMEOUT_SEC", "600") or 600)
WORKER_START_TIMEOUT_SEC = int(os.environ.get("WORKER_START_TIMEOUT_SEC", "120") or 120)
PROCESSED_FILES_PATH = os.environ.get("PROCESSED_FILES_PATH", runtime_path("processed_files.txt"))
ACTIVITY_LOG_PATH = os.environ.get(
    "ACTIVITY_LOG_PATH",
    runtime_path("logs", "activity_turbovec.log") if os.environ.get("VECTOR_BACKEND", "faiss").lower() == "turbovec" else runtime_path("logs", "activity.log"),
)
RESET_DB = False  # DO NOT RESET DB usually

PRIORITY_EXTS = {
    ".xlsx": 0, ".xls": 0,
    ".docx": 1, ".doc": 1,
    ".hwp": 1, ".hwpx": 1, ".txt": 1,
    ".pptx": 2, ".ppt": 2,
    ".pdf": 3,
    ".dwg": 99,
}

STATUS_OK = "ok"
STATUS_NO_CHUNKS = "no_chunks"
STATUS_UNSUPPORTED = "unsupported_extension"
STATUS_EMPTY = "empty_output"
STATUS_ENCRYPTED = "empty_or_encrypted"
STATUS_TIMEOUT = "timeout"
STATUS_NETWORK = "network_error"
STATUS_MISSING_DEP = "missing_dependency"
STATUS_DECODE = "decode_error"
STATUS_PARSE = "parse_error"
STATUS_EMBEDDING = "embedding_error"
STATUS_UNKNOWN = "unknown_error"
STATUS_TEMPORARY = "temporary_file"

class WorkerResponseTimeout(TimeoutError):
    """Raised when the persistent worker stops returning JSON lines."""
    pass

stats_lock = threading.Lock()
extension_stats = defaultdict(Counter)
failure_examples = defaultdict(list)

# Supported extensions and their loaders
LOADERS = {
    ".pdf": PyPDFLoader,
    ".docx": Docx2txtLoader,
    ".doc": LegacyDocLoader, # Now using LibreOffice loader
    ".xlsx": UnstructuredExcelLoader,
    ".xls": UnstructuredExcelLoader,
    ".pptx": PPTXLoader, # Use our custom lighter loader
    ".ppt": LegacyPPTLoader, # Wrap to handle LibreOffice errors gracefully
    ".dwg": DWGLoader,
    ".hwp": HwpLoader,
    ".hwpx": HwpLoader,
    ".txt": TextLoader,
    ".html": UnstructuredHTMLLoader,
    ".htm": UnstructuredHTMLLoader,
}

def get_priority(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    return PRIORITY_EXTS.get(ext, 10)

def classify_error_text(text: str) -> str:
    lower = (text or "").lower()
    if "no module named 'xlrd'" in lower or "missing optional dependency" in lower or "xlrd" in lower:
        return STATUS_MISSING_DEP
    if "unicode" in lower or "codec" in lower or "decode" in lower or "encoding" in lower:
        return STATUS_DECODE
    if "timeout" in lower or "timed out" in lower:
        return STATUS_TIMEOUT
    if "winerror" in lower or "network" in lower or "경로" in lower or "semaphore" in lower:
        return STATUS_NETWORK
    if "encrypted" in lower or "password" in lower:
        return STATUS_ENCRYPTED
    if "unsupported" in lower:
        return STATUS_UNSUPPORTED
    if lower.strip():
        return STATUS_PARSE
    return STATUS_UNKNOWN

def is_temporary_office_file(file_path: str) -> bool:
    return os.path.basename(file_path).startswith("~$")

def record_index_status(file_path: str, status: str, detail: str = ""):
    ext = os.path.splitext(file_path)[1].lower() or "<no_ext>"
    with stats_lock:
        extension_stats[ext][status] += 1
        extension_stats["__total__"][status] += 1
        if detail and len(failure_examples[status]) < 20:
            failure_examples[status].append({
                "file": file_path,
                "detail": detail[:300]
            })
    try:
        from sqlite_index import safe_record_status

        safe_record_status(file_path, status, detail)
    except Exception:
        pass


def record_sidecar_chunks_by_source(vector_docs: List[Dict[str, Any]]) -> None:
    try:
        from sqlite_index import safe_upsert_chunks
    except Exception:
        return
    grouped: Dict[str, List[Dict[str, Any]]] = {}
    for doc in vector_docs or []:
        source = doc.get("source") or doc.get("metadata", {}).get("source") or "Unknown"
        if not source or source == "Unknown":
            continue
        grouped.setdefault(source, []).append({
            "content": doc.get("content", ""),
            "metadata": doc.get("metadata", {}),
        })
    for source, chunks in grouped.items():
        safe_upsert_chunks(source, chunks, status=STATUS_OK)


def checkpoint_sidecar_if_needed(count: int, every: int = 100) -> None:
    if count <= 0 or count % every != 0:
        return
    try:
        from sqlite_index import safe_checkpoint_wal

        safe_checkpoint_wal()
    except Exception:
        pass

def summarize_extension_stats(max_exts: int = 12) -> str:
    with stats_lock:
        rows = []
        for ext, counter in extension_stats.items():
            if ext == "__total__":
                continue
            total = sum(counter.values())
            rows.append((total, ext, dict(counter)))
        rows.sort(reverse=True)
        lines = []
        total_counter = dict(extension_stats.get("__total__", {}))
        if total_counter:
            lines.append(f"total={total_counter}")
        for total, ext, counter in rows[:max_exts]:
            lines.append(f"{ext}: total={total} {counter}")
        return " | ".join(lines) if lines else "no stats yet"

def write_index_report(report_path: str, files_total: int, files_selected: int, started_at: float):
    payload = {
        "started_at": time.ctime(started_at),
        "finished_at": time.ctime(),
        "duration_seconds": round(time.time() - started_at, 2),
        "sample_limit": TEST_LIMIT,
        "sample_mode": INDEX_SAMPLE_MODE,
        "write_enabled": INDEX_SAMPLE_WRITE,
        "files_total": files_total,
        "files_selected": files_selected,
        "by_extension": {},
        "failure_examples": {},
    }
    with stats_lock:
        payload["totals"] = dict(extension_stats.get("__total__", {}))
        for ext, counter in extension_stats.items():
            if ext != "__total__":
                payload["by_extension"][ext] = dict(counter)
        payload["failure_examples"] = {k: v for k, v in failure_examples.items()}
    os.makedirs(os.path.dirname(report_path) or ".", exist_ok=True)
    with open(report_path, "w", encoding="utf-8") as f:
        json.dump(payload, f, ensure_ascii=False, indent=2)
    csv_path = os.path.splitext(report_path)[0] + ".csv"
    with open(csv_path, "w", encoding="utf-8-sig", newline="") as f:
        writer = csv.writer(f)
        statuses = sorted({status for counter in payload["by_extension"].values() for status in counter.keys()})
        writer.writerow(["extension", "total"] + statuses)
        for ext, counter in sorted(payload["by_extension"].items()):
            writer.writerow([ext, sum(counter.values())] + [counter.get(status, 0) for status in statuses])
    print(f"[Report] Wrote index report: {report_path}")

def apply_sample_selection(files: List[str]) -> List[str]:
    if not TEST_LIMIT or len(files) <= TEST_LIMIT:
        return files

    mode = INDEX_SAMPLE_MODE
    if mode == "scan_order":
        print(f"Applying INDEX_SAMPLE_LIMIT={TEST_LIMIT} with scan_order mode.")
        return files[:TEST_LIMIT]

    if mode == "stratified":
        print(f"Applying INDEX_SAMPLE_LIMIT={TEST_LIMIT} with stratified mode.")
        groups = defaultdict(list)
        for file_path in files:
            groups[os.path.splitext(file_path)[1].lower()].append(file_path)
        selected = []
        ordered_exts = sorted(groups.keys(), key=lambda ext: PRIORITY_EXTS.get(ext, 10))
        quota = max(1, TEST_LIMIT // max(1, len(ordered_exts)))
        for ext in ordered_exts:
            selected.extend(groups[ext][:quota])
            if len(selected) >= TEST_LIMIT:
                return selected[:TEST_LIMIT]
        if len(selected) < TEST_LIMIT:
            selected_set = set(selected)
            remaining = [f for f in sorted(files, key=get_priority) if f not in selected_set]
            selected.extend(remaining[:TEST_LIMIT - len(selected)])
        return selected[:TEST_LIMIT]

    print(f"Applying INDEX_SAMPLE_LIMIT={TEST_LIMIT} with priority_first mode.")
    return sorted(files, key=get_priority)[:TEST_LIMIT]

# Lock for tracking file
tracker_lock = threading.Lock()
last_save_time = time.time()

def get_files_from_drives(drives: List[str] | None = None) -> List[str]:
    # Cache file to avoid repeated scanning
    CACHE_FILE = runtime_path("file_list_cache.json")
    
    if os.path.exists(CACHE_FILE):
        try:
            print(f"Loading file list from cache: {CACHE_FILE}...")
            with open(CACHE_FILE, 'r', encoding='utf-8') as f:
                cached_files = json.load(f)
            print(f"Loaded {len(cached_files)} files from cache.")
            return cached_files
        except UnicodeDecodeError:
             print("Cache encoding error. Ignoring cache.")
        except json.JSONDecodeError:
             print("Cache corrupted. Ignoring.")
        except Exception as e:
            print(f"Failed to load cache: {e}. Rescanning...")
    
    all_files = []
    # Always build/load the full file list. Sampling is applied later after
    # processed-file filtering and priority handling so ETA samples are not
    # biased by raw os.walk order.
    target_count = float('inf')
    
    drives = drives or get_search_roots()
    exclude_names = get_exclude_dir_names()
    for drive in drives:
        if len(all_files) >= target_count:
            break
        if not os.path.exists(drive):
            print(f"Warning: Drive {drive} not found. Skipping.")
            continue
        
        print(f"Scanning {drive}...")
        for root, dirs, files in os.walk(drive):
            filter_walk_dirs(dirs, root, exclude_names)
            if len(all_files) >= target_count:
                break
            for file in files:
                ext = os.path.splitext(file)[1].lower()
                if ext in LOADERS:
                    all_files.append(os.path.join(root, file))
                    if len(all_files) >= target_count:
                        break
    
    try:
        print(f"Saving {len(all_files)} files to cache: {CACHE_FILE}...")
        with open(CACHE_FILE, 'w', encoding='utf-8') as f:
            json.dump(all_files, f, ensure_ascii=False, indent=2)
    except Exception as e:
        print(f"Failed to save cache: {e}")

    return all_files

def log_activity(status, file_path):
    """Logs start/end of file processing to tracking file for crash recovery"""
    try:
        with open(ACTIVITY_LOG_PATH, "a", encoding="utf-8") as f:
            f.write(f"{status}|{file_path}\n")
    except:
        pass

def recover_from_crash():
    """Identifies files that were in-flight during a crash and skips them."""
    if TEST_LIMIT and not INDEX_SAMPLE_WRITE:
        print("[Sample] Skipping crash recovery side effects in report-only mode.")
        return
    if not os.path.exists(ACTIVITY_LOG_PATH):
        return

    print("Checking activity log for crash recovery...")
    in_flight = set()
    lines = []
    try:
            with open(ACTIVITY_LOG_PATH, "r", encoding="utf-8", errors="replace") as f:
             lines = f.readlines()
    except Exception as e:
        print(f"Error parse activity log: {e}")
        return

    for line in lines:
        try:
            parts = line.strip().split("|")
            if len(parts) >= 2:
                status, path = parts[0], parts[1]
                if status == "START" or status == "EMBED_START":
                    in_flight.add(path)
                elif status == "END" or status == "EMBED_END":
                    in_flight.discard(path)
        except: pass

    if in_flight:
        print(f"🚨 Found {len(in_flight)} files from crashed session. Marking as SKIPPED.")
        with open(runtime_path("logs", "skipped_crash.txt"), "a", encoding="utf-8") as f:
            for path in in_flight:
                f.write(f"{path}\n")
                print(f"   [Skip] Crashed file: {os.path.basename(path)}")
        
        # Add to processed list so we don't retry them
        append_processed_files(list(in_flight))
        
        # Clear/Rename activity log to reset
        try:
            activity_dir = os.path.dirname(ACTIVITY_LOG_PATH) or "."
            activity_name = os.path.basename(ACTIVITY_LOG_PATH)
            shutil.move(ACTIVITY_LOG_PATH, os.path.join(activity_dir, f"{activity_name}.crash_{int(time.time())}.bak"))
        except:
            pass

def augment_pdf_docs_with_ocr(docs: List[Document], file_path: str) -> List[Document]:
    """Run main-process OCR for PDFs and use OCR text on weak extracted pages."""
    try:
        from ocr_utils import augment_pdf_documents_with_ocr

        augmented_docs = augment_pdf_documents_with_ocr(docs, file_path)
        if augmented_docs:
            return augmented_docs
    except Exception as e:
        print(f"   [OCR] PDF OCR augmentation failed: {os.path.basename(file_path)} | {e}")
    return docs

def load_document(file_path: str) -> List[Document]:
    MAX_RETRIES = 3
    RETRY_DELAY = 5  # seconds
    
    log_activity("START", file_path)
    try:
        ext = os.path.splitext(file_path)[1].lower()
        if is_temporary_office_file(file_path):
            record_index_status(file_path, STATUS_TEMPORARY, "Office temporary lock file")
            return []
        
        # 1. Handle Custom Loaders (PPT/DWG/HWP) in-process (they are safer/custom)
        if ext in [".ppt", ".pptx", ".dwg", ".doc", ".hwp", ".hwpx"]: # Legacy/Custom
            loader_cls = LOADERS.get(ext)
            if loader_cls:
                for attempt in range(MAX_RETRIES):
                    try:
                        loader = loader_cls(file_path)
                        docs = loader.load()
                        if docs:
                            record_index_status(file_path, STATUS_OK)
                        else:
                            record_index_status(file_path, STATUS_NO_CHUNKS, "custom loader returned no documents")
                        return docs
                    except (OSError, IOError) as e:
                        # Network drive error - retry
                        if attempt < MAX_RETRIES - 1:
                            time.sleep(RETRY_DELAY)
                            continue
                        else:
                            send_status_notification(f"⚠️ 네트워크 오류 ({MAX_RETRIES}회 재시도 실패)\n파일: {os.path.basename(file_path)}\n에러: {str(e)[:100]}")
                            record_index_status(file_path, STATUS_NETWORK, str(e))
                            return []
                    except Exception as e:
                        record_index_status(file_path, classify_error_text(str(e)), str(e))
                        return []
            record_index_status(file_path, STATUS_UNSUPPORTED, "no custom loader class")
            return []

        # 2. Handle Standard Loaders via direct import (frozen env 호환)
        for attempt in range(MAX_RETRIES):
            try:
                from worker_loader import load_file

                data = load_file(file_path)

                if isinstance(data, dict) and data.get("__loader_error__"):
                    category = data.get("category", STATUS_PARSE)
                    detail = data.get("detail", "")
                    if category == "empty_or_encrypted":
                        record_index_status(file_path, STATUS_ENCRYPTED, detail)
                    else:
                        record_index_status(file_path, category, detail)
                    return []

                docs = []
                for item in data:
                    docs.append(Document(page_content=item['page_content'], metadata=item['metadata']))
                if ext == ".pdf":
                    docs = augment_pdf_docs_with_ocr(docs, file_path)
                if docs:
                    record_index_status(file_path, STATUS_OK)
                else:
                    record_index_status(file_path, STATUS_NO_CHUNKS, "worker returned zero documents")
                return docs

            except (OSError, IOError) as e:
                # Network drive error - retry
                if attempt < MAX_RETRIES - 1:
                    time.sleep(RETRY_DELAY)
                    continue
                else:
                    send_status_notification(f"⚠️ 네트워크 오류 ({MAX_RETRIES}회 재시도 실패)\n파일: {os.path.basename(file_path)}\n에러: {str(e)[:100]}")
                    record_index_status(file_path, STATUS_NETWORK, str(e))
                    return []
            except Exception as e:
                record_index_status(file_path, classify_error_text(str(e)), str(e))
                return []
        record_index_status(file_path, STATUS_UNKNOWN, "exhausted loader retries")
        return []
    finally:
        log_activity("END", file_path)

def get_existing_sources() -> set:
    if not os.path.exists(VECTOR_STORE_PATH):
        return set()
    print("Checking existing index in ChromaDB (this might take a while)...")
    try:
        from langchain_community.embeddings import HuggingFaceEmbeddings
        embeddings = HuggingFaceEmbeddings(
            model_name="dragonkue/multilingual-e5-small-ko",
            model_kwargs={"device": "cpu"},
            encode_kwargs={"normalize_embeddings": True},
        )
        vectorstore = Chroma(
            persist_directory=VECTOR_STORE_PATH,
            embedding_function=embeddings
        )
        data = vectorstore.get(include=['metadatas'])
        existing = set()
        for meta in data['metadatas']:
            if meta and 'source' in meta:
                existing.add(meta['source'])
        print(f"Found {len(existing)} documents already indexed in DB.")
        return existing
    except Exception as e:
        print(f"Error checking existing index: {e}")
        return set()

def clean_metadata_string(text):
    if not isinstance(text, str):
        return str(text)
    try:
        text.encode('utf-8')
        return text
    except UnicodeEncodeError:
        try:
            original_bytes = text.encode('utf-8', 'surrogateescape')
            return original_bytes.decode('cp949')
        except Exception:
            return text.encode('utf-8', 'replace').decode('utf-8')

# --- New Progress Tracking ---
def load_processed_files() -> set:
    TRACKER_FILE = PROCESSED_FILES_PATH
    if not os.path.exists(TRACKER_FILE):
        return set()
    try:
        with open(TRACKER_FILE, "r", encoding="utf-8") as f:
            return set(line.strip() for line in f)
    except UnicodeDecodeError:
        try:
            print("Warning: Tracker file encoding mismatch. Retrying with CP949...")
            with open(TRACKER_FILE, "r", encoding="cp949") as f:
                return set(line.strip() for line in f)
        except Exception:
             print("Error reading tracker file. Resetting.")
             return set()
    except Exception as e:
        print(f"Error reading tracker file: {e}")
        return set()

def append_processed_files(file_paths: List[str]):
    global last_save_time
    TRACKER_FILE = PROCESSED_FILES_PATH
    with tracker_lock:
        last_save_time = time.time()
        try:
            with open(TRACKER_FILE, "a", encoding="utf-8") as f:
                for p in file_paths:
                    f.write(f"{p}\n")
        except Exception as e:
            print(f"Error updating tracker file: {e}")

# Extended Chroma to allow direct vector insertion with safety batching
class ExtendedChroma(Chroma):
    def add_vectors_direct(self, ids, embeddings, metadatas, documents):
        if not ids:
            return
            
        total = len(ids)
        BATCH_SIZE = 100 # Smaller batch for safety against SQLite limits/Crashes
        
        for i in range(0, total, BATCH_SIZE):
            end = min(i + BATCH_SIZE, total)
            b_ids = ids[i:end]
            b_embeds = embeddings[i:end]
            b_metas = metadatas[i:end]
            b_docs = documents[i:end]
            
            try:
                # Direct call to chromadb collection
                self._collection.add(
                    ids=b_ids,
                    embeddings=b_embeds, 
                    metadatas=b_metas, 
                    documents=b_docs
                )
            except Exception as e:
                print(f"      [Disk] BATCH WRITE FAILED ({i}-{end}): {e}")
                # We do not re-raise to avoid crashing the whole thread?
                # But if we don't raise, we lose data.
                # Losing 100 chunks is better than crashing loop loop.
                # Log detailed error
                with open(runtime_path("logs", "db_write_errors.log"), "a", encoding="utf-8") as f:
                    f.write(f"Batch {i}: {e}\n")

def ingest_data():
    print("Starting ingestion process (GPU Pipeline Mode: Reader -> Embedder -> Writer)...")
    session_started_at = time.time()
    report_path = os.environ.get("INDEX_REPORT_PATH", runtime_path("logs", "index_sample_report.json") if TEST_LIMIT else runtime_path("logs", "index_status_report.json"))
    if TEST_LIMIT:
        print(f"[Sample] INDEX_SAMPLE_LIMIT={TEST_LIMIT}, mode={INDEX_SAMPLE_MODE}, write_enabled={INDEX_SAMPLE_WRITE}")
    
    # 0. Crash Recovery
    recover_from_crash()

    # 1. Get existing files to skip
    existing_sources = load_processed_files()
    print(f"Found {len(existing_sources)} processed files in tracker.")
    
    if not existing_sources and os.path.exists(VECTOR_STORE_PATH):
        print("Tracker empty, checking ChromaDB for existing records (Migration)...")
        chroma_sources = get_existing_sources()
        if chroma_sources:
             print(f"Migrating {len(chroma_sources)} sources from Chroma to tracker...")
             append_processed_files(list(chroma_sources))
             existing_sources = chroma_sources
    
    # 1. Gather files
    files = get_files_from_drives(get_search_roots())
    print(f"Found {len(files)} potential files from cache/scan.")

    seen_candidates = set()
    new_files = []
    duplicate_count = 0
    for file_path in files:
        if file_path in existing_sources:
            continue
        if file_path in seen_candidates:
            duplicate_count += 1
            continue
        seen_candidates.add(file_path)
        new_files.append(file_path)
    if duplicate_count:
        print(f"Excluding {duplicate_count} duplicate file paths from indexing candidates.")
    temp_count = sum(1 for f in new_files if is_temporary_office_file(f))
    if temp_count:
        print(f"Excluding {temp_count} Office temporary lock files (~$*) from indexing candidates.")
        new_files = [f for f in new_files if not is_temporary_office_file(f)]
    
    # --- Priority Sorting ---
    # High Priority: Documents (Fast & Important)
    # Low Priority: DWG (Slow & Heavy)
    print("Sorting files by priority (Docs -> DWG)...")
    new_files.sort(key=get_priority)
    new_files = apply_sample_selection(new_files)
    
    print(f"Skipping {len(files) - len(new_files)} already indexed files.")
    print(f"New files to process: {len(new_files)}")
    if TEST_LIMIT and not INDEX_SAMPLE_WRITE:
        print("[Sample] Report-only mode: vector writes and processed tracker updates are disabled.")
    
    if not new_files:
        print("No new files to ingest.")
        write_index_report(report_path, len(files), 0, session_started_at)
        return

    # --- Queues ---
    # 1. docs_queue: Items are (file_path, List[Document]) or sentinel None
    docs_queue = queue.Queue(maxsize=50) 
    
    # 2. vector_queue: Items are (batch_sources, ids, embeddings, metadatas, texts)
    # 2. vector_queue: Items are (batch_sources, ids, embeddings, metadatas, texts)
    vector_queue = queue.Queue(maxsize=50)

    # --- Process Bridge V4 (Persistent Daemon) ---
    def embedder_bridge_worker(worker_id):
        print(f"   [Process-Bridge] Started. Managing Persistent Processor.")
        
        proc = None
        stdout_queue = None

        def start_stdout_reader(process):
            output_queue = queue.Queue()

            def reader():
                try:
                    for line in process.stdout:
                        output_queue.put(line)
                finally:
                    output_queue.put(None)

            threading.Thread(target=reader, daemon=True).start()
            return output_queue

        def read_processor_line(output_queue, timeout_sec, context):
            try:
                line = output_queue.get(timeout=timeout_sec)
            except queue.Empty:
                raise WorkerResponseTimeout(
                    f"worker_processor stdout timeout after {timeout_sec}s during {context}"
                )

            if line is None:
                raise RuntimeError("Processor died (stdout closed). Crash suspected.")
            return line
        
        def start_processor():
            # Force UTF-8 for subprocess pipes
            env = os.environ.copy()
            env["PYTHONIOENCODING"] = "utf-8"
            
            # Start the persistent worker
            p = subprocess.Popen(
                ["python", "worker_processor.py"],
                stdin=subprocess.PIPE,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE, 
                text=True,
                encoding='utf-8',
                bufsize=1, # Line buffered
                env=env,
                creationflags=subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0,
            )
            output_queue = start_stdout_reader(p)
            # Handshake: Wait for READY
            line = read_processor_line(output_queue, WORKER_START_TIMEOUT_SEC, "startup")
            if "READY" not in line:
                try: p.kill()
                except: pass
                err = p.stderr.read()
                raise RuntimeError(f"Processor init failed: {line} / {err}")
            return p, output_queue

        while True:
            # 1. Ensure Worker is Alive
            if proc is None or proc.poll() is not None:
                print("   [Bridge] Starting Processor Subprocess...")
                try: 
                    proc, stdout_queue = start_processor()
                    print("   [Bridge] Processor Ready.")
                except Exception as e:
                    print(f"   [Bridge] Failed to start processor: {e}. Retrying in 5s...")
                    proc = None
                    stdout_queue = None
                    time.sleep(5)
                    continue

            # 2. Get Job
            try: item = docs_queue.get(timeout=1)
            except queue.Empty: continue
            
            if item is None:
                # Shutdown
                vector_queue.put(None)
                if proc: proc.terminate()
                docs_queue.task_done()
                break
                
            file_path, docs = item
            file_name_only = os.path.basename(file_path)
            
            if not docs:
                vector_queue.put((set(), [], [], [], []))
                docs_queue.task_done(); continue
            
            log_activity("EMBED_START", file_path)
            temp_input = None
            
            try:
                # 3. Validation & Sanitization
                # Pre-Sanitization (to ensure JSON safety)
                valid_docs = [d for d in docs if isinstance(d.page_content, str) and d.page_content.strip()]
                if not valid_docs:
                    record_index_status(file_path, STATUS_NO_CHUNKS, "documents had no usable text")
                    vector_queue.put((set(), [], [], [], []))
                    continue
                
                payload_docs = []
                for d in valid_docs:
                    d.page_content = verify_and_sanitize_content(d.page_content)
                    payload_docs.append({"page_content": d.page_content, "metadata": d.metadata})

                # 4. Create Payload File
                import tempfile
                start_time = time.time()
                with tempfile.NamedTemporaryFile(mode='w', delete=False, encoding='utf-8', suffix='.json') as tmp:
                    json.dump({"docs": payload_docs}, tmp, ensure_ascii=False)
                    temp_input = tmp.name
                
                # 5. Send Command to Worker
                proc.stdin.write(temp_input + "\n")
                proc.stdin.flush()
                
                # 6. Read Response (Blocking)
                # If worker crashes, readline returns empty string immediately (EOF)
                result_line = read_processor_line(
                    stdout_queue,
                    WORKER_FILE_TIMEOUT_SEC,
                    f"processing {file_name_only!r}",
                )
                
                if not result_line:
                    raise RuntimeError("Processor died (EOF). Crash suspected.")
                
                try:
                    result = json.loads(result_line)
                except json.JSONDecodeError:
                    raise RuntimeError(f"Invalid JSON from processor: {result_line}")

                if "error" in result:
                     # Worker caught error but stayed alive, or signalled fatal?
                     # If generic error, we just log and skip this file.
                     raise RuntimeError(f"Processor Error: {result['error']}")
                
                # 7. Success
                ids = result.get("ids", [])
                embeddings = result.get("embeddings", [])
                metadatas = result.get("metadatas", [])
                texts = result.get("texts", [])
                
                duration = time.time() - start_time
                
                print(f"-> [Processor] Processed '{file_name_only}' ({len(texts)} chunks) in {duration:.1f}s")
                batch_sources = {file_path}
                vector_queue.put((batch_sources, ids, embeddings, metadatas, texts))

            except Exception as e:
                try:
                    print(f"🚨 Processor Crash/Error for {file_name_only!r}: {e}")
                except:
                    print("🚨 Processor Crash/Error (Filename encoding error)")
                # If it was a crash (RuntimeError), force restart of process to be safe
                if proc:
                    try: proc.kill()
                    except: pass
                proc = None 
                status = STATUS_TIMEOUT if isinstance(e, WorkerResponseTimeout) else STATUS_EMBEDDING
                record_index_status(file_path, status, str(e))
                
                with open(runtime_path("logs", "skipped_crash.txt"), "a", encoding="utf-8") as f: f.write(f"{file_path}\n")
            
            finally:
                if temp_input and os.path.exists(temp_input):
                    try: os.remove(temp_input)
                    except: pass
                log_activity("EMBED_END", file_path)
                docs_queue.task_done()
        print("   [Process-Bridge] Finished.")

    from langchain_community.embeddings import HuggingFaceEmbeddings # CPU Direct

    # --- Resources ---
    print("Initializing VectorStore and Embeddings (CPU Direct Mode)...")
    # Using HuggingFaceEmbeddings for direct CPU inference
    embeddings_model = HuggingFaceEmbeddings(
        model_name="dragonkue/multilingual-e5-small-ko",
        model_kwargs={'device': 'cpu'},
        encode_kwargs={'normalize_embeddings': True}
    )
    
    # Use ExtendedChroma with telemetry DISABLED via Settings
    import chromadb
    from chromadb.config import Settings
    
    # Use native PersistentClient for proper data persistence
    chroma_settings = Settings(
        anonymized_telemetry=False,
        allow_reset=True
    )
    
    chroma_client = chromadb.PersistentClient(
        path=VECTOR_STORE_PATH,
        settings=chroma_settings
    )
    
    # ExtendedChroma with native client
    vectorstore = ExtendedChroma(
        client=chroma_client,
        collection_name="langchain",
        embedding_function=embeddings_model
    )

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        add_start_index=True,
    )

    # --- Embedder Worker (CPU/APU) - DEAD CODE REMOVED ---
    # (The embedder_worker is effectively replaced by embedder_bridge_worker)

    # --- Writer Worker (VectorStore) ---
    def writer_worker():
        print("   [Writer] Started. VectorStore Mode.")
        if INDEX_SAMPLE_WRITE:
            from faiss_store import add_documents, save_index
        else:
            add_documents = None
            save_index = None
        
        while True:
            item = vector_queue.get()
            
            if item is None: # Sentinel
                if INDEX_SAMPLE_WRITE:
                    save_index()  # 최종 저장
                vector_queue.task_done()
                break
                
            sources, ids, embeddings, metadatas, texts = item
            vector_write_success = False
            vector_docs = []
            
            if ids:
                try:
                    # VectorStore bulk indexing
                    for i, text in enumerate(texts):
                        if embeddings[i] is None or len(embeddings[i]) == 0:
                            continue
                        vector_docs.append({
                            "content": text,
                            "source": metadatas[i].get("source", "Unknown"),
                            "vector": embeddings[i],
                            "metadata": metadatas[i]
                        })
                    
                    if vector_docs and INDEX_SAMPLE_WRITE:
                        add_documents(vector_docs)
                        save_index()
                        record_sidecar_chunks_by_source(vector_docs)
                        vector_write_success = True
                    elif vector_docs:
                        print(f"   [Writer] Sample report-only: skipped vector write for {len(vector_docs)} chunks")
                        vector_write_success = True

                except Exception as e:
                    print(f"   [Writer] Error writing to VectorStore: {e}")
                    for src in sources:
                        record_index_status(src, STATUS_EMBEDDING, str(e))
            
            # Update tracker only after the vector write path succeeds.
            # If add/save fails, leave the file untracked so the next run retries it.
            if sources and ids and vector_write_success:
                unique_to_save = []
                for src in sources:
                     if src not in existing_sources:
                         unique_to_save.append(src)
                         existing_sources.add(src)
                
                if unique_to_save and INDEX_SAMPLE_WRITE:
                    append_processed_files(unique_to_save)
                    checkpoint_sidecar_if_needed(len(existing_sources))
                
            vector_queue.task_done()
        print("   [Writer] Finished.")

    # --- Queue Monitor (Watchdog) ---
    def monitor_queues():
        print("   [Monitor] Started.")
        send_status_notification(f"🚀 인덱싱 시스템 재가동 (최적화 적용)\n전체 대상: {len(files)}개\n남은 작업: {len(new_files)}개")
        
        LOG_FILE = runtime_path("logs", "ingest_status_log.txt")
        with open(LOG_FILE, "a", encoding="utf-8") as f:
            f.write(f"\n--- Session Started: {time.ctime()} ---\n")
        
        last_log_time = 0
        last_notification_time = 0
        last_notification_count = 0
        
        # Stall Detection Variables
        last_stall_check_time = time.time()
        last_stall_check_count = 0
        
        LOG_INTERVAL = 300 # 5 minutes (File Log)
        NOTIFICATION_INTERVAL = 3600 # 1 hour (local notification hook, currently disabled)
        STALL_THRESHOLD = 1800 # 30 minutes
        
        session_start_time = time.time()
        try:
            with open(PROCESSED_FILES_PATH, "r", encoding="utf-8") as f:
                start_done_count = sum(1 for _ in f)
        except:
            start_done_count = 0
            
        sent_30m_report = False
        
        while True:
            time.sleep(10) # check everyone 10s
            d_q = docs_queue.qsize()
            v_q = vector_queue.qsize()
            
            # Estimate processed count from file (Robustness)
            if TEST_LIMIT and not INDEX_SAMPLE_WRITE:
                with stats_lock:
                    done_count = sum(extension_stats.get("__total__", {}).values())
                start_done_count = 0
            else:
                try:
                    with open(PROCESSED_FILES_PATH, "r", encoding="utf-8") as f:
                        done_count = sum(1 for _ in f)
                except:
                    done_count = 0
                
            # Progress Statistics
            session_duration = time.time() - session_start_time
            session_processed = max(0, done_count - start_done_count)
            
            if session_duration > 60: # Calculate after 1 min for stability
                rate_per_min = (session_processed / session_duration) * 60
            else:
                rate_per_min = 0.0
                
            remaining_files = max(0, len(new_files) - session_processed) if TEST_LIMIT else max(0, len(files) - done_count)
            if rate_per_min > 0:
                mins_left = remaining_files / rate_per_min
                eta_hours = int(mins_left // 60)
                eta_mins = int(mins_left % 60)
                eta_message = f"{eta_hours}시간 {eta_mins}분"
            else:
                eta_message = "계산 중..."

            # Progress %
            progress_base = len(new_files) if TEST_LIMIT else len(files)
            progress_done = session_processed if TEST_LIMIT else done_count
            if progress_base > 0:
                progress_pct = (progress_done / progress_base) * 100
            else:
                progress_pct = 0.0

            stats_msg = summarize_extension_stats(5)
            msg = f"[{time.strftime('%H:%M:%S')}] Q: {d_q}/{v_q} | Done: {done_count} ({progress_pct:.2f}%) | Rate: {rate_per_min:.1f}/m | ETA: {eta_message} | Stats: {stats_msg}"
            
            # Console: Always print (Real-time feedback)
            print(f"   [Monitor] {msg}")
            
            current_time = time.time()
            
            # 1. File Logging (Every 5 mins)
            if current_time - last_log_time >= LOG_INTERVAL:
                try:
                    with open(LOG_FILE, "a", encoding="utf-8") as f:
                        f.write(msg + "\n")
                    last_log_time = current_time
                except:
                    pass
            
            # 2. 30 Minute Report (One-time)
            if not sent_30m_report and session_duration >= 1800:
                notification_msg = f"⏱️ 속도 측정 보고 (30분 경과)\n속도: {rate_per_min:.1f}개/분\n남은 작업: {remaining_files}개\n예상 소요: {eta_message}"
                send_status_notification(notification_msg)
                sent_30m_report = True
            
            # 3. Local notification hook (Time Trigger - 1 Hour)
            if current_time - last_notification_time >= NOTIFICATION_INTERVAL:
                notification_msg = f"🔔 정기 생존 신고 (1시간)\n진행률: {progress_pct:.2f}%\n속도: {rate_per_min:.1f}개/분\n예상 소요: {eta_message}\n상태: 안정적"
                send_status_notification(notification_msg)
                last_notification_time = current_time
                
            # 4. Stall Detection (20 Minutes)
            if current_time - last_stall_check_time >= STALL_THRESHOLD: # 30 mins
                if done_count <= last_stall_check_count and done_count > 0:
                      send_status_notification(f"⚠️ 경고: 정체 감지!\n지난 30분간 처리된 파일이 없습니다.\n현재 완료: {done_count}\nETA: {eta_message}")
                
                last_stall_check_time = current_time
                last_stall_check_count = done_count
            
    # --- Start Threads ---
    NUM_EMBEDDERS = 1 # Sequential Mode: One worker to rule them all (avoid locking)
    embedder_threads = []
    
    # Start Monitor
    monitor_thread = threading.Thread(target=monitor_queues, daemon=True)
    monitor_thread.start()
    
    # Use Bridge Worker (Single Process Manager)
    t = threading.Thread(target=embedder_bridge_worker, args=(0,), daemon=True)
    t.start()
    embedder_threads.append(t)
        
    writer_thread = threading.Thread(target=writer_worker, daemon=True)
    writer_thread.start()

    # --- Start Readers (Producers) ---
    MAX_WORKERS = 8 # Ryzen 5 5600G: 6 cores/12 threads, I/O bound
    Bounded_Limit = MAX_WORKERS * 8
    total_processed_count = 0
    total_new = len(new_files)
    
    print(f"Starting Readers with {MAX_WORKERS} worker (Sequential Subprocess Mode)...")
    
    file_iterator = iter(new_files)
    
    with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
        futures = {executor.submit(load_document, f): f for f in itertools.islice(file_iterator, Bounded_Limit)}
        
        while futures:
            done, _ = wait(futures, return_when=FIRST_COMPLETED)
            
            for future in done:
                file_path = futures.pop(future)
                total_processed_count += 1
                
                if total_processed_count % 50 == 0:
                     print(f"   [Reader] {total_processed_count}/{total_new} files read... ", end="\r")
                
                try:
                    docs = future.result()
                    # Push key info to queue
                    docs_queue.put((file_path, docs))
                except Exception as e:
                    record_index_status(file_path, classify_error_text(str(e)), str(e))
                    print(f"   [Reader] Error reading {os.path.basename(file_path)}: {e}")
                
                next_file = next(file_iterator, None)
                if next_file:
                    futures[executor.submit(load_document, next_file)] = next_file

    print("\nAll files read via Readers.")
    
    # Poison Pills
    for _ in range(NUM_EMBEDDERS):
        docs_queue.put(None)
    for t in embedder_threads:
        t.join()
        
    vector_queue.put(None)
    writer_thread.join()
    write_index_report(report_path, len(files), len(new_files), session_started_at)
    
    print("\nIngestion complete!")

if __name__ == "__main__":
    try:
        ingest_data()
    except KeyboardInterrupt:
        print("\nProcess stopped by user.")
    except Exception as e:
        import traceback
        with open(runtime_path("logs", "crash_log.txt"), "w", encoding="utf-8") as f:
            f.write(traceback.format_exc())
        print(f"\nCRITICAL ERROR: {e}")
        send_windows_notification("RAG Ingestion Error", f"Process crashed: {e}")
        raise e
